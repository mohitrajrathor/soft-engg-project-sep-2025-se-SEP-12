"""
PPTX Export Service for Slide Decks
Generates styled PowerPoint presentations with themes, colors, and embedded charts.
"""

from io import BytesIO
import base64
import re
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image


class SlideTheme:
    """Defines a visual theme for slides."""
    def __init__(
        self,
        name: str,
        bg_color: RGBColor,
        title_color: RGBColor,
        text_color: RGBColor,
        accent_color: RGBColor,
        font_title: str = "Calibri",
        font_body: str = "Calibri"
    ):
        self.name = name
        self.bg_color = bg_color
        self.title_color = title_color
        self.text_color = text_color
        self.accent_color = accent_color
        self.font_title = font_title
        self.font_body = font_body


# Predefined themes
THEMES = {
    "professional": SlideTheme(
        name="Professional",
        bg_color=RGBColor(255, 255, 255),  # White
        title_color=RGBColor(31, 78, 120),  # Navy Blue
        text_color=RGBColor(64, 64, 64),    # Dark Gray
        accent_color=RGBColor(68, 114, 196), # Blue
        font_title="Calibri Light",
        font_body="Calibri"
    ),
    "modern": SlideTheme(
        name="Modern",
        bg_color=RGBColor(248, 249, 250),   # Light Gray
        title_color=RGBColor(13, 110, 253), # Bright Blue
        text_color=RGBColor(33, 37, 41),    # Almost Black
        accent_color=RGBColor(111, 66, 193), # Purple
        font_title="Arial",
        font_body="Arial"
    ),
    "colorful": SlideTheme(
        name="Colorful",
        bg_color=RGBColor(255, 248, 240),   # Peach
        title_color=RGBColor(220, 53, 69),  # Red
        text_color=RGBColor(52, 58, 64),    # Dark Gray
        accent_color=RGBColor(253, 126, 20), # Orange
        font_title="Georgia",
        font_body="Verdana"
    ),
    "dark": SlideTheme(
        name="Dark",
        bg_color=RGBColor(33, 37, 41),      # Dark Gray
        title_color=RGBColor(255, 193, 7),  # Yellow
        text_color=RGBColor(248, 249, 250), # White
        accent_color=RGBColor(32, 201, 151), # Teal
        font_title="Segoe UI",
        font_body="Segoe UI"
    ),
    "minimalist": SlideTheme(
        name="Minimalist",
        bg_color=RGBColor(250, 250, 250),   # Off-White
        title_color=RGBColor(0, 0, 0),      # Black
        text_color=RGBColor(96, 96, 96),    # Gray
        accent_color=RGBColor(144, 144, 144), # Light Gray
        font_title="Arial",
        font_body="Arial"
    )
}


class PPTXExportService:
    """Service for exporting slide decks to PPTX format with styling."""
    
    def __init__(self, theme_name: str = "professional"):
        """Initialize with a theme."""
        self.theme = THEMES.get(theme_name.lower(), THEMES["professional"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)  # 16:9 aspect ratio
    
    def _parse_markdown_to_runs(self, text_frame, markdown_text: str):
        """
        Parse markdown text and apply formatting.
        Supports: **bold**, *italic*, `code`, and basic color syntax.
        """
        # Split by lines for paragraphs
        lines = markdown_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if it's a bullet point
            is_bullet = line.startswith('- ') or line.startswith('* ')
            if is_bullet:
                line = line[2:].strip()
            
            p = text_frame.add_paragraph()
            if is_bullet:
                p.level = 0
            
            # Parse inline formatting
            self._parse_inline_formatting(p, line)
            
            # Apply theme styling
            for run in p.runs:
                run.font.name = self.theme.font_body
                run.font.size = Pt(16)
                # Set color only if not already set
                try:
                    if run.font.color.rgb is None:
                        run.font.color.rgb = self.theme.text_color
                except AttributeError:
                    # Color type doesn't support rgb, set it directly
                    run.font.color.rgb = self.theme.text_color
    
    def _parse_inline_formatting(self, paragraph, text: str):
        """Parse and apply inline formatting (bold, italic, code)."""
        # Pattern to match **bold**, *italic*, `code`
        pattern = r'(\*\*.*?\*\*|\*.*?\*|`.*?`)'
        parts = re.split(pattern, text)
        
        for part in parts:
            if not part:
                continue
            
            run = paragraph.add_run()
            
            # Bold
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            # Italic
            elif part.startswith('*') and part.endswith('*'):
                run.text = part[1:-1]
                run.font.italic = True
            # Code (monospace)
            elif part.startswith('`') and part.endswith('`'):
                run.text = part[1:-1]
                run.font.name = 'Courier New'
                run.font.color.rgb = self.theme.accent_color
            # Plain text
            else:
                run.text = part
    
    def _add_title_slide(self, title: str, subtitle: Optional[str] = None):
        """Add a title slide with theme styling."""
        slide_layout = self.prs.slide_layouts[0]  # Title layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.bg_color
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.theme.font_title
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = self.theme.title_color
        
        # Subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            for paragraph in subtitle_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.theme.font_body
                    run.font.size = Pt(24)
                    run.font.color.rgb = self.theme.text_color
    
    def _add_content_slide(self, title: str, content: str, graph_image: Optional[str] = None):
        """Add a content slide with optional chart image."""
        # Choose layout based on whether there's an image
        if graph_image:
            slide_layout = self.prs.slide_layouts[5]  # Blank layout for custom positioning
        else:
            slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.bg_color
        
        # Add title
        if graph_image:
            # Manual title for blank layout
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.theme.font_title
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = self.theme.title_color
        else:
            title_shape = slide.shapes.title
            title_shape.text = title
            title_frame = title_shape.text_frame
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.theme.font_title
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = self.theme.title_color
        
        # Add content
        if graph_image:
            # Two-column layout: text left, image right
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.2), Inches(4.5), Inches(4)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            self._parse_markdown_to_runs(text_frame, content)
            
            # Add image
            try:
                image_data = self._decode_base64_image(graph_image)
                slide.shapes.add_picture(
                    image_data,
                    Inches(5.2), Inches(1.2),
                    width=Inches(4.3), height=Inches(3.8)
                )
            except Exception as e:
                print(f"Warning: Failed to add image to slide: {e}")
        else:
            # Single column content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            self._parse_markdown_to_runs(text_frame, content)
    
    def _decode_base64_image(self, base64_string: str) -> BytesIO:
        """Decode base64 image string (with or without data URI prefix)."""
        # Remove data URI prefix if present
        if ',' in base64_string:
            base64_string = base64_string.split(',', 1)[1]
        
        image_data = base64.b64decode(base64_string)
        return BytesIO(image_data)
    
    def generate_presentation(
        self,
        deck_title: str,
        deck_description: Optional[str],
        slides: List[dict]
    ) -> BytesIO:
        """
        Generate a complete PowerPoint presentation.
        
        Args:
            deck_title: Title of the deck
            deck_description: Optional description
            slides: List of slide dictionaries with 'title', 'content', 'graph_image' keys
        
        Returns:
            BytesIO object containing the PPTX file
        """
        # Add title slide
        self._add_title_slide(deck_title, deck_description)
        
        # Add content slides
        for slide_data in slides:
            self._add_content_slide(
                title=slide_data.get('title', 'Untitled'),
                content=slide_data.get('content', ''),
                graph_image=slide_data.get('graph_image')
            )
        
        # Save to BytesIO
        pptx_bytes = BytesIO()
        self.prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return pptx_bytes


def export_to_pptx(
    deck_title: str,
    deck_description: Optional[str],
    slides: List[dict],
    theme: str = "professional"
) -> BytesIO:
    """
    Main function to export a slide deck to PPTX format.
    
    Args:
        deck_title: Title of the presentation
        deck_description: Optional subtitle/description
        slides: List of slide dictionaries
        theme: Theme name (professional, modern, colorful, dark, minimalist)
    
    Returns:
        BytesIO containing the PPTX file
    """
    service = PPTXExportService(theme_name=theme)
    return service.generate_presentation(deck_title, deck_description, slides)
